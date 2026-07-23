from __future__ import annotations

import json
import os
import shutil
from pathlib import Path
from tempfile import TemporaryDirectory

import pikepdf
import pypdfium2

from ..config import AppPaths, ensure_app_dirs
from ..release_config import ReleaseConfig
from .artifact_service import ArtifactService
from .knowledge_content_store import KnowledgeContentStore
from .knowledge_derivation_service import KnowledgeDerivationService
from .knowledge_import_service import KnowledgeImportService
from .knowledge_import_worker import read_worker_result
from .knowledge_service import KnowledgeService
from .knowledge_vector_store import KnowledgeVectorRecord, LanceKnowledgeVectorStore
from .paddle_ocr_service import (
    PaddleOcrDeploymentService,
    PaddleOcrState,
    ReleasePaddleOcrBundleSource,
)
from .storage import StorageBootstrapService
from .storage.layout import knowledge_import_result_path


def run_knowledge_packaged_smoke(paths: AppPaths) -> None:
    """Exercise Knowledge native/data boundaries from the running executable."""

    # Keep heavy document runtimes operation-scoped. Importing Docling mutates the
    # process-wide ElementTree namespace registry, which must not affect unrelated
    # SVG generation merely because this smoke module was discovered.
    from docling_core.types.doc import DocItemLabel, DoclingDocument
    from docx import Document
    from pptx import Presentation

    # OCR model paths are passed to a native Windows process. Keep this isolated
    # smoke topology deliberately short so it tests product behavior rather than
    # inheriting an arbitrarily deep CI/pytest temporary root.
    with TemporaryDirectory(prefix="xk-", dir=paths.temp) as temporary:
        root = Path(temporary)
        smoke_paths = ensure_app_dirs(
            AppPaths(
                home=root,
                config=root / "c",
                logs=root / "l",
                cache=root / "k",
                state=root / "s",
                temp=root / "t",
                artifacts=root / "a",
                resources=paths.resources,
            )
        )
        pdf_path = root / "knowledge-smoke.pdf"
        _write_simple_pdf(pdf_path, "Knowledge packaged smoke")

        with pikepdf.Pdf.open(pdf_path) as document:
            if len(document.pages) != 1:
                raise RuntimeError("pikepdf packaged Knowledge smoke failed.")

        pdfium_document = pypdfium2.PdfDocument(pdf_path)
        try:
            page = pdfium_document[0]
            try:
                image = page.render(scale=0.25).to_pil()
                if image.width < 1 or image.height < 1:
                    raise RuntimeError("PDFium packaged Knowledge render failed.")
            finally:
                page.close()
        finally:
            pdfium_document.close()

        seed_document = DoclingDocument(name="knowledge-packaged-smoke")
        seed_document.add_text(
            label=DocItemLabel.PARAGRAPH,
            text="Knowledge packaged smoke",
        )
        docling_document = seed_document
        if "Knowledge packaged smoke" not in docling_document.export_to_text():
            raise RuntimeError("Docling packaged Knowledge IR parse failed.")

        worker_docx = root / "knowledge-worker-smoke.docx"
        worker_document = Document()
        worker_document.add_paragraph("Knowledge frozen worker smoke")
        worker_document.save(worker_docx)
        worker_pptx = root / "knowledge-worker-smoke.pptx"
        worker_presentation = Presentation()
        slide = worker_presentation.slides.add_slide(
            worker_presentation.slide_layouts[1]
        )
        slide.shapes.title.text = "Knowledge presentation smoke"
        slide.placeholders[1].text = "Knowledge frozen presentation worker smoke"
        worker_presentation.save(worker_pptx)
        content_store = KnowledgeContentStore(smoke_paths)
        stored = content_store.write_canonical_bundle(
            envelope={"canonical_generation_id": "packaged-smoke"},
            docling_document=docling_document.export_to_dict(),
        )
        reopened = content_store.read_canonical_bundle(
            stored.relative_path,
            expected_envelope_sha256=stored.envelope_sha256,
            expected_content_ir_sha256=stored.content_ir_sha256,
        )
        if reopened.envelope.get("canonical_generation_id") != "packaged-smoke":
            raise RuntimeError("Zstandard canonical packaged Knowledge smoke failed.")

        storage = StorageBootstrapService().initialize(smoke_paths)
        derivation = KnowledgeDerivationService(
            paths=smoke_paths,
            session_factory=storage.session_factory,
            start_worker=False,
        )
        knowledge = KnowledgeService(storage.session_factory)
        importer = KnowledgeImportService(
            paths=smoke_paths,
            session_factory=storage.session_factory,
            artifact_service=ArtifactService(storage.session_factory),
            canonical_ready_notifier=derivation.enqueue_generation,
        )
        try:
            imported_by_format = {}
            for source_format, source in (
                ("docx", worker_docx),
                ("pptx", worker_pptx),
            ):
                imported = importer.import_file(source, timeout=180)
                import_worker_result = read_worker_result(
                    knowledge_import_result_path(smoke_paths, imported.import_id)
                )
                if import_worker_result.worker_pid == os.getpid():
                    raise RuntimeError("Knowledge import did not use a spawned worker.")
                if (
                    import_worker_result.status != "succeeded"
                    or import_worker_result.failure_stage is not None
                    or import_worker_result.diagnostic_code is not None
                ):
                    raise RuntimeError(
                        f"Spawned {source_format.upper()} import worker result is invalid."
                    )
                imported_by_format[source_format] = imported
            presentation_import = imported_by_format["pptx"]
            derivation_view = derivation.status_for_import(presentation_import.import_id)
            if derivation_view is None:
                raise RuntimeError("Spawned PPTX derivation was not queued.")
            derived = derivation.derive_now(derivation_view.job_id)
            matches = knowledge.lookup("presentation worker smoke", top_k=5)
            if (
                not derived.retrieval_ready
                or not matches
                or "presentation worker smoke" not in matches[0].quote.casefold()
            ):
                raise RuntimeError("Spawned PPTX import did not reach Knowledge lookup.")
        finally:
            importer.shutdown()
            derivation.shutdown()
            storage.engine.dispose()

        vector_store = LanceKnowledgeVectorStore(smoke_paths)
        relative_path = vector_store.write_generation(
            generation_id="packaged-smoke",
            records=(
                KnowledgeVectorRecord("unit-a", (1.0, 0.0, 0.0)),
                KnowledgeVectorRecord("unit-b", (0.0, 1.0, 0.0)),
            ),
            dimensions=3,
            corpus_fingerprint="corpus-smoke",
            profile_fingerprint="profile-smoke",
        )
        if vector_store.search(relative_path, query_vector=(0.9, 0.1, 0.0), limit=1) != ["unit-a"]:
            raise RuntimeError("LanceDB packaged Knowledge search failed.")

        ocr_archive_value = os.environ.get("XENIX_KNOWLEDGE_OCR_SMOKE_ARCHIVE", "").strip()
        ocr_deployment = PaddleOcrDeploymentService(smoke_paths)
        if ocr_deployment.status_snapshot().state is not PaddleOcrState.NOT_INSTALLED:
            raise RuntimeError("Fresh native OCR deployment state is invalid.")
        catalog = ocr_deployment.catalog
        if catalog is not None and (
            catalog.protocol_version != 2
            or not catalog.runtime_id.startswith("paddle-inference-")
            or not catalog.artifact_name.endswith(".zip")
        ):
            raise RuntimeError("Packaged native OCR catalog is invalid.")
        native_ocr_activated = False
        native_ocr_retrieval = False
        if ocr_archive_value:
            if catalog is None:
                raise RuntimeError("Packaged native OCR catalog is missing.")
            ocr_archive = Path(ocr_archive_value).resolve()
            if not ocr_archive.is_file() or ocr_archive.name != catalog.artifact_name:
                raise RuntimeError("Packaged native OCR smoke archive is invalid.")
            cached_archive = (
                smoke_paths.cache / "knowledge-ocr" / "downloads" / catalog.artifact_name
            )
            cached_archive.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(ocr_archive, cached_archive)
            ocr_deployment = PaddleOcrDeploymentService(
                smoke_paths,
                bundle_source=ReleasePaddleOcrBundleSource(
                    catalog,
                    ReleaseConfig(
                        releases_oss_public_url="https://knowledge-ocr-smoke.invalid"
                    ),
                ),
            )
            phases: list[str] = []
            installed = ocr_deployment.install(phases.append)
            if installed.state is not PaddleOcrState.READY or phases[-1:] != ["ready"]:
                raise RuntimeError("Packaged native OCR activation failed.")
            runtime = ocr_deployment.open_runtime()
            if not runtime.executable_path.is_file():
                raise RuntimeError("Packaged native OCR executable is unavailable after activation.")
            native_ocr_activated = True

            ocr_image_value = os.environ.get(
                "XENIX_KNOWLEDGE_OCR_SMOKE_IMAGE",
                "",
            ).strip()
            ocr_image = Path(ocr_image_value).resolve() if ocr_image_value else None
            if ocr_image is None or not ocr_image.is_file():
                raise RuntimeError("Packaged native OCR smoke image is unavailable.")
            image_header = ocr_image.read_bytes()[:8]
            image_suffix = ".jpg" if image_header.startswith(b"\xff\xd8\xff") else ".png"
            import_image = root / f"native-ocr-retrieval{image_suffix}"
            shutil.copy2(ocr_image, import_image)
            ocr_storage = StorageBootstrapService().initialize(smoke_paths)
            knowledge = KnowledgeService(ocr_storage.session_factory)
            derivation = KnowledgeDerivationService(
                paths=smoke_paths,
                session_factory=ocr_storage.session_factory,
                start_worker=False,
            )
            importer = KnowledgeImportService(
                paths=smoke_paths,
                session_factory=ocr_storage.session_factory,
                artifact_service=ArtifactService(ocr_storage.session_factory),
                canonical_ready_notifier=derivation.enqueue_generation,
            )
            try:
                imported = importer.import_file(import_image, timeout=180)
                derivation_view = derivation.status_for_import(imported.import_id)
                if derivation_view is None:
                    raise RuntimeError("Native OCR derivation was not queued.")
                derived = derivation.derive_now(derivation_view.job_id)
                if not derived.retrieval_ready:
                    raise RuntimeError("Native OCR retrieval projection is unavailable.")
                matches = knowledge.lookup("BOARDING", top_k=5)
                if not matches or "BOARDING" not in matches[0].quote:
                    raise RuntimeError("Native OCR text did not reach Knowledge lookup.")
                worker_result = read_worker_result(
                    knowledge_import_result_path(smoke_paths, imported.import_id)
                )
                recorded_runtime = worker_result.pipeline.get("ocr", {}).get("runtime")
                if not isinstance(recorded_runtime, dict) or (
                    recorded_runtime.get("generation_id") != runtime.generation_id
                ):
                    raise RuntimeError("Native OCR runtime provenance was not recorded.")
                native_ocr_retrieval = True
            finally:
                importer.shutdown()
                derivation.shutdown()
                ocr_storage.engine.dispose()

    marker = paths.state / "knowledge-smoke.json"
    marker.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "docling_ir": True,
                "pdfium_render": True,
                "pikepdf": True,
                "canonical_zstd": True,
                "import_worker_spawn": True,
                "spawned_docx_import": True,
                "spawned_pptx_import": True,
                "lancedb": True,
                "paddle_native_deployment": True,
                "paddle_native_activation": native_ocr_activated,
                "paddle_native_retrieval": native_ocr_retrieval,
            },
            sort_keys=True,
        ),
        encoding="utf-8",
    )


def _write_simple_pdf(path: Path, text: str) -> None:
    content = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    objects = (
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        f"<< /Length {len(content)} >>\nstream\n".encode("ascii") + content + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    )
    payload = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, body in enumerate(objects, start=1):
        offsets.append(len(payload))
        payload.extend(f"{index} 0 obj\n".encode("ascii"))
        payload.extend(body)
        payload.extend(b"\nendobj\n")
    xref = len(payload)
    payload.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    payload.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        payload.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    payload.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode("ascii")
    )
    path.write_bytes(payload)
